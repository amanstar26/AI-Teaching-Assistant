from pptx import Presentation
from pptx.util import Pt
import re

def create_ppt(
    content,
    filename="chapter_presentation.pptx"
    ):

    prs = Presentation()

    slide_blocks = re.split(
        r"SLIDE\s+\d+",
        content,
        flags=re.IGNORECASE
    )

    slide_blocks = [
        block.strip()
        for block in slide_blocks
        if block.strip()
    ]

    if not slide_blocks:

        slide = prs.slides.add_slide(
            prs.slide_layouts[1]
        )

        slide.shapes.title.text = (
            "Presentation"
        )

        slide.placeholders[1].text = (
            content[:1000]
        )

        prs.save(filename)

        return filename

    for block in slide_blocks:

        title = "Untitled Slide"

        body_points = []

        lines = [
            line.strip()
            for line in block.splitlines()
            if line.strip()
        ]

        for line in lines:

            title_match = re.match(
                r"TITLE\s*:\s*(.*)",
                line,
                flags=re.IGNORECASE
            )

            if title_match:

                title = title_match.group(1).strip()

                continue

            if (
                line.upper().startswith("POINTS")
                or
                line.upper().startswith("CONTENT")
            ):
                continue

            line = (
                line.replace("•", "")
                .replace("-", "")
                .strip()
            )

            if line:

                body_points.append(line)

        slide = prs.slides.add_slide(
            prs.slide_layouts[1]
        )

        slide.shapes.title.text = title

        for paragraph in (
            slide.shapes.title.text_frame.paragraphs
        ):

            paragraph.font.size = Pt(24)

            paragraph.font.bold = True

        text_frame = (
            slide.placeholders[1]
            .text_frame
        )

        text_frame.clear()

        for i, point in enumerate(body_points):

            if i == 0:

                text_frame.text = point

                text_frame.paragraphs[0].font.size = Pt(18)

            else:

                p = text_frame.add_paragraph()

                p.text = point

                p.level = 0

                p.font.size = Pt(18)

    prs.save(filename)

    return filename
